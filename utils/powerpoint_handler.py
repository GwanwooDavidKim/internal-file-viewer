# -*- coding: utf-8 -*-
"""
PowerPoint 파일 처리 모듈 (PowerPoint File Handler)

python-pptx를 사용하여 PowerPoint 파일에서 텍스트와 슬라이드 정보를 추출합니다.
PDF 변환을 통한 안전한 미리보기 기능을 제공합니다.

핵심 개선사항:
- 사용자의 PowerPoint 작업에 영향 없음 (PDF 변환 방식)
- 원본 파일 락 없음
- "원본 열기" 기능 완벽 작동
- 빠르고 안정적인 렌더링
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from .ppt_to_pdf_converter import get_converter
from .com_powerpoint_converter import get_com_converter
from .aspose_powerpoint_converter import get_aspose_converter
from .pdf_handler import PdfHandler
import logging
import time

# PIL을 안전하게 import
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("Warning: Pillow not installed. PowerPoint image preview will not be available.")

logger = logging.getLogger(__name__)


class PowerPointHandler:
    """
    PowerPoint 파일 처리를 위한 클래스입니다.
    
    주요 기능:
    - PowerPoint 슬라이드 텍스트 추출
    - 안전한 PDF 변환을 통한 미리보기
    - 슬라이드 구조 분석
    - 프레젠테이션 메타데이터 조회
    - 사용자 작업 방해 없음 (완전 격리)
    """
    
    def __init__(self):
        """PowerPointHandler 인스턴스를 초기화합니다."""
        # PDF 변환기와 PDF 핸들러 초기화 (먼저 생성)
        # 우선순위: Aspose (평가판, 워터마크 허용) → LibreOffice (무료) → 텍스트만 추출
        self.aspose_converter = get_aspose_converter()
        self.pdf_converter = get_converter()  # LibreOffice 변환기
        self.pdf_handler = PdfHandler()
        
        # 사용할 변환기 결정 (우선순위 적용)
        if self.aspose_converter.is_available():
            self.active_converter = self.aspose_converter
            self.converter_type = "Aspose.Slides (평가판)"
            self.supported_extensions = ['.ppt', '.pptx']  # Aspose는 모든 PowerPoint 형식 지원
            print("   [시작] Aspose.Slides 방식 사용 (평가판 - 워터마크 허용)")
            print("   [파일] 지원 형식: .ppt, .pptx")
            print("   [안전] Microsoft Office 설치 불필요")
            print("   [평가판] 워터마크 허용 - 사용자 간섭 없음")
        else:
            self.active_converter = self.pdf_converter
            self.converter_type = "LibreOffice" 
            self.supported_extensions = ['.pptx']  # LibreOffice는 .pptx만 안정적
            print("   [캐시] LibreOffice 방식 사용 (호환성)")
            print("   [파일] 지원 형식: .pptx")
        
        # 현재 연결된 파일 경로 (호환성을 위해)
        self.current_file_path = None
        
        print("[처리중] PowerPointHandler 초기화 - 안전한 PDF 변환 방식 사용")
        print("   [완료] 사용자 PowerPoint 작업에 영향 없음")
        print("   [완료] 원본 파일 락 없음") 
        print("   [완료] '원본 열기' 기능 완벽 작동")
        print(f"   [변환기] 활성 변환기: {self.converter_type}")
        
        # 변환기별 특징 안내
        if self.converter_type.startswith("Aspose"):
            print("   [작업] Aspose 장점: 사용자 간섭 없음 + LibreOffice보다 빠름 (워터마크 포함)")
        else:
            print("   [무료] 무료 솔루션 사용")
    
    def open_persistent_connection(self, file_path: str) -> bool:
        """
        호환성을 위한 메소드 - PDF 변환 방식에서는 지속 연결이 불필요
        
        Args:
            file_path (str): PowerPoint 파일 경로
            
        Returns:
            bool: 항상 True (PDF 변환 방식은 항상 사용 가능)
        """
        self.current_file_path = file_path  # 현재 파일 경로 저장 (render_slide_fast용)
        logger.info(f"[처리중] PPT → PDF 방식으로 연결: {os.path.basename(file_path)}")
        logger.info("   [완료] 지속 연결 불필요 - 즉시 렌더링 가능")
        return True
    
    def close_persistent_connection(self):
        """
        호환성을 위한 메소드 - PDF 변환 방식에서는 정리할 연결이 없음
        """
        self.current_file_path = None  # 현재 파일 경로 초기화
        logger.info("[처리중] PPT → PDF 방식 정리 완료")
        logger.info("   [완료] 사용자 PowerPoint에 영향 없이 안전하게 종료")
    
    def is_connected(self) -> bool:
        """
        호환성을 위한 메소드 - 연결 상태 확인
        
        Returns:
            bool: 파일이 연결되어 있는지 여부
        """
        return self.current_file_path is not None
    
    def render_slide_fast(self, slide_number: int, width: int = 800, height: int = 600) -> Optional['Image.Image']:
        """
        호환성을 위한 메소드 - PDF 변환 방식에서는 빠른/일반 렌더링 구분이 없음
        
        Args:
            slide_number (int): 슬라이드 번호 (0부터 시작)
            width (int): 이미지 너비
            height (int): 이미지 높이
            
        Returns:
            Optional[Image.Image]: 렌더링된 이미지
        """
        if not self.current_file_path:
            logger.error("[오류] render_slide_fast 호출 전에 open_persistent_connection이 필요합니다")
            return None
            
        logger.info(f"[시작] 빠른 렌더링 (PDF 방식): 슬라이드 {slide_number + 1}")
        # PDF 변환 방식은 항상 빠르므로 기본 렌더링 메소드와 동일
        return self.render_slide_to_image(self.current_file_path, slide_number, width, height)
    
    def can_handle(self, file_path: str) -> bool:
        """
        파일이 이 핸들러가 처리할 수 있는 형식인지 확인합니다.
        
        Args:
            file_path (str): 파일 경로
            
        Returns:
            bool: 처리 가능 여부
        """
        return any(file_path.lower().endswith(ext) for ext in self.supported_extensions)
    
    def get_slide_count(self, file_path: str) -> int:
        """
        PowerPoint의 총 슬라이드 수를 반환합니다.
        
        Args:
            file_path (str): PowerPoint 파일 경로
            
        Returns:
            int: 슬라이드 수 (오류 시 0)
        """
        # .ppt 파일이나 Aspose가 활성인 경우 Aspose 사용
        if (file_path.lower().endswith('.ppt') or 
            self.converter_type.startswith("Aspose")):
            if hasattr(self.active_converter, 'get_slide_count'):
                try:
                    count = self.active_converter.get_slide_count(file_path)
                    if count > 0:
                        return count
                except Exception as e:
                    logger.error(f"Aspose 슬라이드 수 확인 오류: {e}")
        
        # .pptx 파일은 python-pptx로 직접 처리
        try:
            prs = Presentation(file_path)
            return len(prs.slides)
        except Exception as e:
            logger.error(f"슬라이드 수 확인 오류: {e}")
            return 0
    
    def render_slide_to_image(self, file_path: str, slide_number: int, width: int = 800, height: int = 600) -> Optional['Image.Image']:
        """
        PPT → PDF → 이미지 방식으로 안전하게 슬라이드를 렌더링합니다.
        
        이 방식의 장점:
        - 사용자의 PowerPoint 작업에 영향 없음
        - 원본 파일이 잠기지 않음  
        - "원본 열기" 기능 완벽 작동
        - 빠르고 안정적인 렌더링
        
        Args:
            file_path (str): PowerPoint 파일 경로
            slide_number (int): 슬라이드 번호 (0부터 시작)
            width (int): 이미지 너비 (미사용 - PDF 기본 해상도)
            height (int): 이미지 높이 (미사용 - PDF 기본 해상도)
            
        Returns:
            Optional[Image.Image]: 생성된 이미지 (실패 시 None)
        """
        if not PIL_AVAILABLE:
            logger.warning("PIL이 없어서 이미지 렌더링 불가")
            return None
        
        try:
            logger.info(f"[처리중] PPT → PDF → 이미지 렌더링 시작: {os.path.basename(file_path)}, 슬라이드 {slide_number + 1}")
            
            # 1단계: PPT를 PDF로 변환 (캐시 활용) - 활성 변환기 사용
            start_time = time.time()
            pdf_path = self.active_converter.convert_to_pdf(file_path)
            conversion_time = time.time() - start_time
            if not pdf_path:
                logger.error("[오류] PPT → PDF 변환 실패")
                return None
            
            logger.info(f"[완료] PDF 변환 완료: {os.path.basename(pdf_path)}")
            
            # 2단계: PDF에서 해당 페이지를 이미지로 렌더링
            image = self.pdf_handler.render_page_to_image(
                pdf_path, 
                page_num=slide_number,
                zoom=1.5  # 고품질을 위한 150% 확대
            )
            
            if image:
                logger.info(f"[완료] 슬라이드 {slide_number + 1} 렌더링 완료! ({self.converter_type} 변환: {conversion_time:.1f}초)")
                return image
            else:
                logger.error(f"[오류] PDF 페이지 {slide_number} 렌더링 실패")
                return None
                
        except Exception as e:
            logger.error(f"[오류] 슬라이드 렌더링 오류: {e}")
            return None
    
    def extract_text_from_slide(self, file_path: str, slide_number: int) -> Dict[str, Any]:
        """
        지정된 슬라이드에서 텍스트를 추출합니다.
        
        Args:
            file_path (str): PowerPoint 파일 경로
            slide_number (int): 슬라이드 번호 (0부터 시작)
            
        Returns:
            Dict[str, Any]: 슬라이드 정보
        """
        try:
            prs = Presentation(file_path)
            
            if slide_number >= len(prs.slides) or slide_number < 0:
                return {'error': '잘못된 슬라이드 번호'}
            
            slide = prs.slides[slide_number]
            
            # 슬라이드 제목 추출
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
            
            # 텍스트 내용 추출
            text_content = []
            bullet_points = []
            
            for shape in slide.shapes:
                # 텍스트가 있는 shape만 처리
                if (hasattr(shape, "text") and hasattr(shape, "text_frame") and 
                    hasattr(shape, 'text') and shape.text and shape.text.strip()):
                    # 제목이 아닌 경우에만 추가
                    if shape != slide.shapes.title:
                        text_content.append(shape.text)
                        
                        # 텍스트 프레임이 있는 경우 단락별로 분석
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    bullet_points.append({
                                        'text': paragraph.text,
                                        'level': paragraph.level,
                                    })
            
            # 이미지 및 기타 객체 카운트
            image_count = 0
            chart_count = 0
            table_count = 0
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    chart_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_count += 1
            
            return {
                'slide_number': slide_number + 1,
                'title': title or '[제목 없음]',
                'text_content': text_content,
                'bullet_points': bullet_points,
                'full_text': "\n".join(text_content),
                'image_count': image_count,
                'chart_count': chart_count,
                'table_count': table_count,
                'total_shapes': len(slide.shapes),
            }
            
        except Exception as e:
            logger.error(f"슬라이드 텍스트 추출 오류: {e}")
            return {'error': f"슬라이드 텍스트 추출 오류: {e}"}
    
    def extract_text_by_slides(self, file_path: str, max_slides: int = None) -> List[Dict[str, Any]]:
        """
        전체 프레젠테이션에서 슬라이드별로 텍스트를 추출합니다.
        
        Args:
            file_path (str): PowerPoint 파일 경로
            max_slides (int, optional): 최대 슬라이드 수 제한 (None이면 모든 슬라이드)
            
        Returns:
            List[Dict[str, Any]]: 슬라이드별 텍스트 정보 [{"page_num": 1, "content": "..."}, ...]
        """
        # .ppt 파일은 python-pptx로 직접 읽을 수 없으므로 PDF에서 텍스트 추출
        if file_path.lower().endswith('.ppt'):
            try:
                logger.info(f"[처리중] .ppt 파일 텍스트 추출: PDF 변환 방식 사용")
                pdf_path = self.active_converter.convert_to_pdf(file_path)
                if pdf_path:
                    return self.pdf_handler.extract_text_by_pages(pdf_path, max_pages=max_slides)
                else:
                    return [{"page_num": 1, "content": f".ppt 파일 텍스트 추출 실패: {os.path.basename(file_path)}"}]
            except Exception as e:
                logger.error(f".ppt 텍스트 추출 오류: {e}")
                return [{"page_num": 1, "content": f".ppt 파일 텍스트 추출 오류: {e}"}]
        
        # .pptx 파일은 python-pptx로 직접 추출
        try:
            prs = Presentation(file_path)
            slides_data = []
            
            # 슬라이드 수 제한 적용
            slides_to_process = prs.slides
            if max_slides is not None:
                slides_to_process = list(prs.slides)[:max_slides]
            
            for i, slide in enumerate(slides_to_process):
                slide_text_parts = []
                
                # 슬라이드 제목
                title = ""
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()
                
                # 슬라이드 내용
                for shape in slide.shapes:
                    if (hasattr(shape, "text") and hasattr(shape, "text_frame") and 
                        hasattr(shape, 'text') and shape.text and shape.text.strip()):
                        if shape != slide.shapes.title:
                            slide_text_parts.append(shape.text)
                
                # 제목과 내용 결합
                content_parts = []
                if title:
                    content_parts.append(title)
                if slide_text_parts:
                    content_parts.extend(slide_text_parts)
                
                slides_data.append({
                    "page_num": i + 1,
                    "content": "\n".join(content_parts) if content_parts else "[텍스트 없음]"
                })
            
            return slides_data
            
        except Exception as e:
            logger.error(f"PowerPoint 텍스트 추출 오류: {e}")
            return [{"page_num": 1, "content": f"PowerPoint 텍스트 추출 오류: {e}"}]
    
    def extract_text(self, file_path: str, max_slides: int = None) -> str:
        """
        전체 프레젠테이션에서 텍스트를 추출합니다.
        (검색 인덱싱용)
        
        Args:
            file_path (str): PowerPoint 파일 경로
            max_slides (int, optional): 최대 슬라이드 수 제한 (None이면 모든 슬라이드)
            
        Returns:
            str: 추출된 전체 텍스트
        """
        try:
            slides_data = self.extract_text_by_slides(file_path, max_slides)
            all_text = []
            
            for slide_data in slides_data:
                slide_num = slide_data["page_num"]
                content = slide_data["content"]
                all_text.append(f"=== 슬라이드 {slide_num} ===\n{content}")
            
            return "\n\n".join(all_text)
            
        except Exception as e:
            logger.error(f"PowerPoint 텍스트 추출 오류: {e}")
            return f"PowerPoint 텍스트 추출 오류: {e}"
    
    def extract_all_text(self, file_path: str, max_slides: int = None) -> str:
        """extract_text의 별칭 (호환성을 위해)"""
        return self.extract_text(file_path, max_slides)
    
    def get_presentation_info(self, file_path: str) -> Dict[str, Any]:
        """
        프레젠테이션의 상세 정보를 반환합니다.
        
        Args:
            file_path (str): PowerPoint 파일 경로
            
        Returns:
            Dict[str, Any]: 프레젠테이션 정보
        """
        try:
            if not os.path.exists(file_path):
                return {'error': '파일을 찾을 수 없습니다'}
            
            # .ppt 파일은 python-pptx로 읽을 수 없으므로 PDF 기반 정보 추출
            if file_path.lower().endswith('.ppt'):
                return self._get_ppt_info_via_pdf(file_path)
            
            # .pptx 파일은 python-pptx로 직접 처리
            prs = Presentation(file_path)
            
            # 기본 정보
            file_size = os.path.getsize(file_path)
            slide_count = len(prs.slides)
            
            # 프레젠테이션 메타데이터
            core_props = prs.core_properties
            
            # 슬라이드 크기 정보
            slide_width = prs.slide_width or 9144000  # 기본값 (10인치)
            slide_height = prs.slide_height or 6858000  # 기본값 (7.5인치)
            
            # 각 슬라이드의 요약 정보
            slides_summary = []
            total_images = 0
            total_charts = 0
            total_tables = 0
            
            for i, slide in enumerate(prs.slides):
                # 슬라이드 제목
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text
                
                # 객체 카운트
                images = charts = tables = 0
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        images += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        charts += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        tables += 1
                
                total_images += images
                total_charts += charts
                total_tables += tables
                
                # 텍스트 블록 수
                text_shapes = sum(1 for shape in slide.shapes 
                                if (hasattr(shape, "text") and hasattr(shape, "text_frame") and 
                                    hasattr(shape, 'text') and shape.text and shape.text.strip()))
                
                slides_summary.append({
                    'slide_number': i + 1,
                    'title': title or f'[슬라이드 {i + 1}]',
                    'text_shapes': text_shapes,
                    'images': images,
                    'charts': charts,
                    'tables': tables,
                    'total_shapes': len(slide.shapes),
                })
            
            # PDF 변환 가능 여부 확인 (활성 변환기에서 직접)
            conversion_info = self.active_converter.get_cache_info()
            conversion_available = conversion_info.get('converter_available', 
                                                      conversion_info.get('libreoffice_available', False))
            
            info = {
                'filename': os.path.basename(file_path),
                'file_size': file_size,
                'file_size_mb': round(file_size / (1024 * 1024), 2),
                'slide_count': slide_count,
                'slide_width_inches': round(slide_width / 914400, 2),
                'slide_height_inches': round(slide_height / 914400, 2),
                'total_images': total_images,
                'total_charts': total_charts,
                'total_tables': total_tables,
                'slides_summary': slides_summary,
                'conversion_available': conversion_available,
                'converter_type': self.converter_type,
                'cache_info': conversion_info,
                'metadata': {
                    'title': getattr(core_props, 'title', None),
                    'subject': getattr(core_props, 'subject', None),
                    'author': getattr(core_props, 'author', None),
                    'created': getattr(core_props, 'created', None),
                    'last_modified_by': getattr(core_props, 'last_modified_by', None),
                    'modified': getattr(core_props, 'modified', None),
                }
            }
            
            return info
            
        except Exception as e:
            logger.error(f"프레젠테이션 정보 조회 오류: {e}")
            return {'error': f'프레젠테이션 정보 조회 오류: {e}'}
    
    def _get_ppt_info_via_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        .ppt 파일의 정보를 PDF 변환을 통해 추출합니다.
        
        Args:
            file_path (str): .ppt 파일 경로
            
        Returns:
            Dict[str, Any]: 프레젠테이션 정보
        """
        try:
            logger.info(f"[처리중] .ppt 파일 정보 추출: PDF 변환 방식 사용")
            
            # 기본 파일 정보
            file_size = os.path.getsize(file_path)
            
            # PDF로 변환하여 슬라이드 수 확인
            pdf_path = self.active_converter.convert_to_pdf(file_path)
            slide_count = 0
            if pdf_path:
                # PDF 핸들러로 페이지 수 확인 (= 슬라이드 수)
                slide_count = self.pdf_handler.get_page_count(pdf_path)
            
            # 변환 정보
            conversion_info = self.active_converter.get_cache_info()
            conversion_available = conversion_info.get('converter_available', 
                                                      conversion_info.get('libreoffice_available', False))
            
            return {
                'filename': os.path.basename(file_path),
                'file_size': file_size,
                'file_size_mb': round(file_size / (1024 * 1024), 2),
                'slide_count': slide_count,
                'file_type': '.ppt (Legacy PowerPoint)',
                'conversion_available': conversion_available,
                'converter_type': self.converter_type,
                'cache_info': conversion_info,
                'note': '.ppt 파일은 PDF 변환을 통해서만 미리보기 가능합니다.',
                'slides_summary': [{'slide_number': i+1, 'title': f'슬라이드 {i+1}', 'note': 'PDF 변환 방식'} 
                                  for i in range(slide_count)],
                'metadata': {
                    'title': None,
                    'subject': None, 
                    'author': None,
                    'created': None,
                    'last_modified_by': None,
                    'modified': None,
                }
            }
            
        except Exception as e:
            logger.error(f".ppt 파일 정보 추출 오류: {e}")
            return {'error': f'.ppt 파일 정보 추출 오류: {e}'}
    
    def search_in_presentation(self, file_path: str, search_term: str, max_results: int = 10) -> List[Dict[str, Any]]:
        """
        프레젠테이션에서 텍스트를 검색합니다.
        
        Args:
            file_path (str): PowerPoint 파일 경로
            search_term (str): 검색어
            max_results (int): 최대 결과 수
            
        Returns:
            List[Dict[str, Any]]: 검색 결과
        """
        try:
            prs = Presentation(file_path)
            results = []
            search_term_lower = search_term.lower()
            
            for i, slide in enumerate(prs.slides):
                # 슬라이드 제목 검색
                if slide.shapes.title and slide.shapes.title.text:
                    title_text = slide.shapes.title.text
                    if search_term_lower in title_text.lower():
                        results.append({
                            'slide_number': i + 1,
                            'location': '제목',
                            'type': 'title',
                            'text': title_text,
                            'context': title_text,
                        })
                        
                        if len(results) >= max_results:
                            return results
                
                # 슬라이드 내용 검색
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text") and hasattr(shape, "text_frame") and shape.text.strip():
                        if shape != slide.shapes.title:
                            text = shape.text
                            if search_term_lower in text.lower():
                                # 컨텍스트 생성 (검색어 주변 텍스트)
                                context_start = max(0, text.lower().find(search_term_lower) - 50)
                                context_end = min(len(text), text.lower().find(search_term_lower) + len(search_term) + 50)
                                context = text[context_start:context_end]
                                
                                if context_start > 0:
                                    context = "..." + context
                                if context_end < len(text):
                                    context = context + "..."
                                
                                results.append({
                                    'slide_number': i + 1,
                                    'location': f'텍스트 블록 {shape_idx + 1}',
                                    'type': 'content',
                                    'text': text,
                                    'context': context,
                                })
                                
                                if len(results) >= max_results:
                                    return results
            
            return results
            
        except Exception as e:
            logger.error(f"프레젠테이션 검색 오류: {e}")
            return [{'error': f"프레젠테이션 검색 오류: {e}"}]